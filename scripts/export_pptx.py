#!/usr/bin/env python3
"""
export_pptx.py — Export HTML Slide Deck to PowerPoint (.pptx)
══════════════════════════════════════════════════════════════════════════════
Converts a single-file HTML presentation back to an editable .pptx file.
Completes the two-way workflow:

  PPT  →  extract_pptx.py  →  HTML          (already existed)
  HTML →  export_pptx.py   →  PPT           (this script — new direction)

Extraction pipeline:
  1. parse_html.py parses the HTML into a JSON outline
  2. export_pptx.py reads the outline and builds the PPTX via python-pptx

Usage:
    python3 scripts/export_pptx.py <input.html> [options]

Examples:
    python3 scripts/export_pptx.py presentation.html
    python3 scripts/export_pptx.py deck.html --output my_deck.pptx
    python3 scripts/export_pptx.py deck.html --theme dark     # Dark background
    python3 scripts/export_pptx.py deck.html --no-notes       # Skip speaker notes
    python3 scripts/export_pptx.py deck.html --verbose

Options:
    --output, -o    Output .pptx path (default: <input_stem>.pptx)
    --theme         Slide theme: light | dark | auto (default: auto — inferred from HTML)
    --no-notes      Skip speaker notes (default: notes ARE included)
    --verbose, -v   Print per-slide progress

Requirements:
    pip3 install python-pptx beautifulsoup4
══════════════════════════════════════════════════════════════════════════════
"""

import argparse
import json
import re
import sys
from pathlib import Path


# ── Dependency checks ─────────────────────────────────────────────────────────

def check_deps():
    missing = []
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        missing.append("python-pptx")
    try:
        from bs4 import BeautifulSoup  # noqa: F401
    except ImportError:
        missing.append("beautifulsoup4")
    if missing:
        print(f"❌ Missing dependencies: {', '.join(missing)}")
        print(f"   Run: pip3 install {' '.join(missing)}")
        sys.exit(1)


# ── Color utilities ───────────────────────────────────────────────────────────

def hex_to_rgb(hex_color):
    """Convert #rrggbb to (r, g, b) tuple."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    if len(hex_color) != 6:
        return (30, 30, 46)  # dark fallback
    try:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except ValueError:
        return (30, 30, 46)


def detect_theme_from_html(html_text):
    """Infer light/dark theme from CSS variable definitions."""
    dark_bg_patterns = [
        r'--bg\s*:\s*#[01][0-9a-f]{5}',
        r'--clr-bg\s*:\s*#[01][0-9a-f]{5}',       # healthcare/finance/education CSS var name
        r'background-color\s*:\s*#[01][0-9a-f]{5}',
        r'class=["\'][^"\']*theme-pitch-deck',
        r'class=["\'][^"\']*theme-tech-talk',
        r'class=["\'][^"\']*theme-hhart-red',
        r'class=["\'][^"\']*theme-product-launch',
        r'class=["\'][^"\']*theme-pash-orange',
        r'class=["\'][^"\']*theme-finance',         # Bloomberg Dark (new)
        # Bloomberg Dark specific: charcoal #111827
        r'--clr-bg\s*:\s*#111827',
        r'--bg\s*:\s*#111827',
    ]
    for pat in dark_bg_patterns:
        if re.search(pat, html_text, re.I):
            return "dark"
    return "light"


def extract_css_color(html_text, var_name):
    """Extract a CSS custom property value like --primary: #c9a84c"""
    pat = re.compile(rf'{re.escape(var_name)}\s*:\s*(#[0-9a-fA-F]{{3,6}})')
    m = pat.search(html_text)
    return m.group(1) if m else None


# ── HTML → structured outline ─────────────────────────────────────────────────

def html_to_outline(html_path):
    """
    Reuse parse_html logic inline to avoid subprocess dependency.
    Returns a list of slide dicts with keys: type, title, subtitle, items,
    body, stats, notes, left, right, quote, author.
    """
    from bs4 import BeautifulSoup

    html = html_path.read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "html.parser")

    def clean(text):
        return re.sub(r'\s+', ' ', (text or "").strip())

    slide_els = soup.select(".slide") or soup.select("section")
    slides = []

    for idx, el in enumerate(slide_els):
        # Speaker notes
        notes = ""
        notes_el = el.select_one(".notes, aside.notes")
        if notes_el:
            notes = clean(notes_el.get_text())
            notes_el.decompose()
        if el.get("data-notes"):
            notes = notes or clean(el["data-notes"])

        # Basic fields
        classes = " ".join(el.get("class", []))
        has_list = bool(el.find("ul") or el.find("ol"))
        has_bq   = bool(el.find("blockquote") or el.select_one(".quote-text"))
        has_stat = bool(el.select(".stat-card, .stat-item, .metric, .kpi"))
        is_title = idx == 0 or bool(el.find("h1"))
        is_end   = "slide-end" in classes or "slide-thanks" in classes
        two_col  = bool(el.select(".col, .column, .two-col"))

        slide = {"notes": notes}

        if is_end:
            slide["type"] = "end"
            h = el.select_one("h2, h3, .slide-title")
            slide["title"] = clean(h.get_text()) if h else "Thank You"
            sub = el.select_one(".subtitle")
            if sub: slide["subtitle"] = clean(sub.get_text())

        elif is_title:
            slide["type"] = "title"
            h1 = el.select_one("h1, .title-text")
            slide["title"] = clean(h1.get_text()) if h1 else ""
            sub = el.select_one(".subtitle, h2")
            if sub: slide["subtitle"] = clean(sub.get_text())
            attr = el.select_one(".attr, .author, .byline")
            if attr: slide["author"] = clean(attr.get_text())

        elif has_bq:
            slide["type"] = "quote"
            q = el.select_one("blockquote, .quote-text")
            slide["quote"] = clean(q.get_text()) if q else ""
            attr = el.select_one(".attr, cite, .attribution")
            if attr: slide["author"] = clean(attr.get_text())

        elif has_stat:
            slide["type"] = "stats"
            h = el.select_one("h2, h3, .slide-title")
            if h: slide["title"] = clean(h.get_text())
            stats = []
            for card in el.select(".stat-card, .stat-item, .metric, .kpi"):
                val = card.select_one(".stat-value, .value, .number, strong")
                lbl = card.select_one(".stat-label, .label, span:last-child")
                if val:
                    s = {"value": clean(val.get_text())}
                    if lbl and lbl != val:
                        s["label"] = clean(lbl.get_text())
                    stats.append(s)
            slide["stats"] = stats

        elif two_col:
            slide["type"] = "two-col"
            h = el.select_one("h2, h3, .slide-title")
            if h: slide["title"] = clean(h.get_text())
            cols = el.select(".col, .column, .left, .right, [class*='col-']")
            if len(cols) >= 2:
                def col_data(col_el):
                    t = col_el.select_one("h3, h4, strong, .col-title")
                    b = col_el.select_one("p, .col-body")
                    return {
                        "title": clean(t.get_text()) if t else "",
                        "body":  clean(b.get_text()) if b else clean(col_el.get_text()),
                    }
                slide["left"]  = col_data(cols[0])
                slide["right"] = col_data(cols[1])

        elif has_list:
            slide["type"] = "bullets"
            h = el.select_one("h2, h3, .slide-title")
            if h: slide["title"] = clean(h.get_text())
            slide["items"] = [clean(li.get_text()) for li in el.select("li") if clean(li.get_text())]

        else:
            slide["type"] = "text"
            h = el.select_one("h2, h3, .slide-title")
            if h: slide["title"] = clean(h.get_text())
            sub = el.select_one(".subtitle")
            if sub: slide["subtitle"] = clean(sub.get_text())
            p = el.select_one("p, .body")
            if p: slide["body"] = clean(p.get_text())

        slides.append(slide)

    return html, slides


# ── PPTX builder ──────────────────────────────────────────────────────────────

def build_pptx(slides, html_text, theme, include_notes, verbose):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Widescreen 16:9
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme colors — search multiple CSS variable names used across all 11 templates
    is_dark = (theme == "dark")
    bg_color_hex = (extract_css_color(html_text, "--bg")
                    or extract_css_color(html_text, "--clr-bg")
                    or ("#0d0d1a" if is_dark else "#fafaf8"))
    accent_hex   = (extract_css_color(html_text, "--primary")
                    or extract_css_color(html_text, "--accent")
                    or extract_css_color(html_text, "--clr-accent")
                    or "#6366f1")
    text_hex     = "#ffffff" if is_dark else "#1a1a2e"

    bg_rgb     = RGBColor(*hex_to_rgb(bg_color_hex))
    accent_rgb = RGBColor(*hex_to_rgb(accent_hex))
    text_rgb   = RGBColor(*hex_to_rgb(text_hex))
    muted_rgb  = RGBColor(160, 160, 180) if is_dark else RGBColor(100, 100, 120)

    W = prs.slide_width
    H = prs.slide_height

    def add_bg(slide):
        """Fill slide background with theme color."""
        from pptx.util import Pt
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

    def add_text_box(slide, text, left, top, width, height,
                     font_size=24, bold=False, color=None, align=PP_ALIGN.LEFT,
                     word_wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color or text_rgb
        return txBox

    def add_accent_bar(slide, width_frac=0.08):
        """Left vertical accent bar."""
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Emu(0), Emu(0), int(W * width_frac), H
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_rgb
        bar.line.fill.background()

    # ── Slide type renderers ──────────────────────────────────────────────────

    def render_title(slide, data):
        add_bg(slide)
        # Accent bar
        bar = slide.shapes.add_shape(1, Emu(0), Emu(0), int(W * 0.006), H)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent_rgb; bar.line.fill.background()
        # Title
        title = data.get("title", "")
        add_text_box(slide, title,
                     int(W * 0.08), int(H * 0.3), int(W * 0.84), int(H * 0.25),
                     font_size=44, bold=True, align=PP_ALIGN.LEFT)
        # Subtitle
        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle,
                         int(W * 0.08), int(H * 0.58), int(W * 0.7), int(H * 0.15),
                         font_size=20, color=muted_rgb)
        # Author
        author = data.get("author", "")
        if author:
            add_text_box(slide, author,
                         int(W * 0.08), int(H * 0.76), int(W * 0.5), int(H * 0.1),
                         font_size=14, color=muted_rgb)

    def render_bullets(slide, data):
        add_bg(slide)
        title = data.get("title", "")
        if title:
            add_text_box(slide, title,
                         int(W * 0.07), int(H * 0.06), int(W * 0.86), int(H * 0.15),
                         font_size=32, bold=True)
        # Accent underline
        line = slide.shapes.add_shape(1,
            int(W * 0.07), int(H * 0.21), int(W * 0.12), int(H * 0.004))
        line.fill.solid(); line.fill.fore_color.rgb = accent_rgb; line.line.fill.background()

        items = data.get("items", [])
        y_start = int(H * 0.25)
        row_h   = int(H * 0.09)
        for i, item in enumerate(items[:7]):
            y = y_start + i * row_h
            # Bullet dot
            dot = slide.shapes.add_shape(1,
                int(W * 0.07), y + int(row_h * 0.35),
                int(W * 0.012), int(H * 0.012))
            dot.fill.solid(); dot.fill.fore_color.rgb = accent_rgb; dot.line.fill.background()
            add_text_box(slide, item,
                         int(W * 0.1), y, int(W * 0.84), row_h,
                         font_size=18)

    def render_stats(slide, data):
        add_bg(slide)
        title = data.get("title", "")
        if title:
            add_text_box(slide, title,
                         int(W * 0.07), int(H * 0.06), int(W * 0.86), int(H * 0.14),
                         font_size=28, bold=True)
        stats = data.get("stats", [])
        n = min(len(stats), 4)
        if n == 0:
            return
        card_w = int(W * 0.8 / n)
        x_start = int(W * 0.1)
        for i, stat in enumerate(stats[:4]):
            x = x_start + i * (card_w + int(W * 0.02))
            y = int(H * 0.3)
            # Card bg
            card = slide.shapes.add_shape(1, x, y, card_w, int(H * 0.45))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(
                min(hex_to_rgb(bg_color_hex)[0] + 20, 255),
                min(hex_to_rgb(bg_color_hex)[1] + 20, 255),
                min(hex_to_rgb(bg_color_hex)[2] + 20, 255),
            )
            card.line.fill.background()
            # Value
            add_text_box(slide, stat.get("value", "—"),
                         x + int(card_w * 0.1), y + int(H * 0.06),
                         int(card_w * 0.8), int(H * 0.2),
                         font_size=36, bold=True, color=accent_rgb, align=PP_ALIGN.CENTER)
            # Label
            label = stat.get("label", "")
            if label:
                add_text_box(slide, label,
                             x + int(card_w * 0.05), y + int(H * 0.28),
                             int(card_w * 0.9), int(H * 0.12),
                             font_size=14, color=muted_rgb, align=PP_ALIGN.CENTER)

    def render_quote(slide, data):
        add_bg(slide)
        # Decorative quote mark
        add_text_box(slide, "\u201c",
                     int(W * 0.07), int(H * 0.05), int(W * 0.15), int(H * 0.25),
                     font_size=120, color=accent_rgb)
        quote = data.get("quote", "")
        add_text_box(slide, quote,
                     int(W * 0.12), int(H * 0.25), int(W * 0.76), int(H * 0.45),
                     font_size=24, align=PP_ALIGN.LEFT)
        author = data.get("author", "")
        if author:
            add_text_box(slide, f"— {author}",
                         int(W * 0.12), int(H * 0.72), int(W * 0.76), int(H * 0.1),
                         font_size=16, color=muted_rgb)

    def render_two_col(slide, data):
        add_bg(slide)
        title = data.get("title", "")
        if title:
            add_text_box(slide, title,
                         int(W * 0.07), int(H * 0.06), int(W * 0.86), int(H * 0.13),
                         font_size=28, bold=True)
        # Divider
        div = slide.shapes.add_shape(1,
            int(W * 0.5), int(H * 0.22), int(W * 0.002), int(H * 0.65))
        div.fill.solid(); div.fill.fore_color.rgb = muted_rgb; div.line.fill.background()

        left  = data.get("left", {})
        right = data.get("right", {})
        for col_data, x_off in [(left, 0.06), (right, 0.52)]:
            ct = col_data.get("title", "")
            cb = col_data.get("body", "")
            if ct:
                add_text_box(slide, ct,
                             int(W * x_off), int(H * 0.22), int(W * 0.4), int(H * 0.1),
                             font_size=20, bold=True, color=accent_rgb)
            if cb:
                add_text_box(slide, cb,
                             int(W * x_off), int(H * 0.34), int(W * 0.4), int(H * 0.5),
                             font_size=16)

    def render_text(slide, data):
        add_bg(slide)
        title = data.get("title", "")
        if title:
            add_text_box(slide, title,
                         int(W * 0.07), int(H * 0.12), int(W * 0.86), int(H * 0.16),
                         font_size=32, bold=True)
        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle,
                         int(W * 0.07), int(H * 0.3), int(W * 0.86), int(H * 0.1),
                         font_size=20, color=muted_rgb)
        body = data.get("body", "")
        if body:
            add_text_box(slide, body,
                         int(W * 0.07), int(H * 0.42), int(W * 0.86), int(H * 0.45),
                         font_size=18)

    def render_end(slide, data):
        add_bg(slide)
        # Full-width accent block at top
        top = slide.shapes.add_shape(1, Emu(0), Emu(0), W, int(H * 0.008))
        top.fill.solid(); top.fill.fore_color.rgb = accent_rgb; top.line.fill.background()
        title = data.get("title", "Thank You")
        add_text_box(slide, title,
                     int(W * 0.1), int(H * 0.3), int(W * 0.8), int(H * 0.25),
                     font_size=48, bold=True, align=PP_ALIGN.CENTER)
        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle,
                         int(W * 0.1), int(H * 0.58), int(W * 0.8), int(H * 0.12),
                         font_size=20, color=muted_rgb, align=PP_ALIGN.CENTER)

    # ── Blank layout index 6 (most templates have a blank layout) ─────────────
    blank_layout = prs.slide_layouts[6]

    RENDERERS = {
        "title":    render_title,
        "bullets":  render_bullets,
        "stats":    render_stats,
        "quote":    render_quote,
        "two-col":  render_two_col,
        "features": render_bullets,   # Render features as bullets fallback
        "text":     render_text,
        "end":      render_end,
    }

    for i, slide_data in enumerate(slides):
        sl = prs.slides.add_slide(blank_layout)
        slide_type = slide_data.get("type", "text")
        renderer = RENDERERS.get(slide_type, render_text)
        renderer(sl, slide_data)

        if include_notes and slide_data.get("notes"):
            notes_slide = sl.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data["notes"]

        if verbose:
            print(f"   Slide {i+1:02d}: [{slide_type}] {slide_data.get('title', '')[:40]}")

    return prs


# ── Main ──────────────────────────────────────────────────────────────────────

def export_pptx(input_path, output_path=None, theme="auto",
                include_notes=True, verbose=False):
    input_path = Path(input_path)
    if not input_path.exists():
        print(f"❌ File not found: {input_path}")
        sys.exit(1)

    if verbose:
        print(f"📥 Parsing: {input_path.name}")

    html_text, slides = html_to_outline(input_path)

    if not slides:
        print("⚠️  No slides found. Is this a frontend-slides presentation?")
        sys.exit(1)

    # Theme detection
    if theme == "auto":
        theme = detect_theme_from_html(html_text)
        if verbose:
            print(f"   Detected theme: {theme}")

    if verbose:
        print(f"   Found {len(slides)} slides — building PPTX…")

    prs = build_pptx(slides, html_text, theme, include_notes, verbose)

    if output_path is None:
        output_path = input_path.with_suffix(".pptx")
    else:
        output_path = Path(output_path)

    prs.save(str(output_path))

    size_kb = round(output_path.stat().st_size / 1024)
    print(f"✅ Exported → {output_path}  ({len(slides)} slides, {size_kb} KB)")
    print(f"   Open with: open '{output_path}'")


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Export an HTML slide presentation to a PowerPoint .pptx file",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python3 scripts/export_pptx.py presentation.html
  python3 scripts/export_pptx.py deck.html --output my_deck.pptx --theme dark
  python3 scripts/export_pptx.py deck.html --notes --verbose
        """
    )
    parser.add_argument("input", help="Path to the HTML presentation file")
    parser.add_argument("--output", "-o", help="Output .pptx path (default: <input_stem>.pptx)")
    parser.add_argument("--theme", choices=["auto", "light", "dark"], default="auto",
                        help="Slide theme: auto (inferred), light, or dark")
    parser.add_argument("--notes", action="store_true", default=True,
                        help="Include speaker notes — default ON (use --no-notes to disable)")
    parser.add_argument("--no-notes", action="store_false", dest="notes",
                        help="Skip speaker notes in the exported PPTX")
    parser.add_argument("--verbose", "-v", action="store_true", help="Print per-slide progress")
    args = parser.parse_args()

    check_deps()
    export_pptx(
        args.input,
        output_path=args.output,
        theme=args.theme,
        include_notes=args.notes,
        verbose=args.verbose,
    )


if __name__ == "__main__":
    main()
