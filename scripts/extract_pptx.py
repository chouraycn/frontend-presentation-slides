#!/usr/bin/env python3
"""
extract_pptx.py — Extract content from PowerPoint files for frontend-slides skill.

Usage:
    python3 extract_pptx.py <input.pptx> --output <output-dir>

Output:
    <output-dir>/
        slides.json       — Structured slide data (includes tables, SmartArt,
                            text styles, improved layout detection)
        images/           — Extracted images (slide_N_img_M.png)
        summary.txt       — Human-readable content summary

Requirements:
    pip3 install python-pptx Pillow

Changelog v2:
    - Table extraction: captures headers + rows as structured JSON
    - SmartArt extraction: reads underlying XML text nodes
    - Richer text styles: italic, color, alignment, underline
    - Improved layout detection: 8 layout types, scoring-based
    - Position-aware title detection: uses left/top geometry, not just shape name
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path


def check_dependencies():
    """Check if required packages are available."""
    missing = []
    try:
        import pptx
    except ImportError:
        missing.append("python-pptx")
    try:
        from PIL import Image
    except ImportError:
        missing.append("Pillow")

    if missing:
        print(f"❌ Missing required packages: {', '.join(missing)}")
        print(f"   Run: pip3 install {' '.join(missing)}")
        sys.exit(1)


# ── Text style extraction ────────────────────────────────────────────────────

def _run_style(run):
    """Return a style dict for a single text run."""
    style = {}
    try:
        if run.font.bold is not None:
            style["bold"] = bool(run.font.bold)
        if run.font.italic is not None:
            style["italic"] = bool(run.font.italic)
        if run.font.underline is not None:
            style["underline"] = bool(run.font.underline)
        if run.font.size:
            style["size"] = int(run.font.size.pt)
        # Color — RGBColor or None
        try:
            color = run.font.color.rgb if run.font.color and run.font.color.type else None
            if color:
                style["color"] = f"#{color}"
        except Exception:
            pass
    except Exception:
        pass
    return style


def _para_style(para):
    """Return the dominant style for a paragraph (first non-None wins)."""
    style = {}
    try:
        # Alignment
        if para.alignment:
            from pptx.enum.text import PP_ALIGN
            align_map = {
                PP_ALIGN.LEFT: "left",
                PP_ALIGN.CENTER: "center",
                PP_ALIGN.RIGHT: "right",
                PP_ALIGN.JUSTIFY: "justify",
            }
            style["align"] = align_map.get(para.alignment, "left")
    except Exception:
        pass

    # Merge run styles — first defined value wins
    for run in para.runs:
        rs = _run_style(run)
        for k, v in rs.items():
            if k not in style:
                style[k] = v
    return style


def extract_text_from_shape(shape):
    """Extract text paragraphs from any text-bearing shape."""
    paragraphs = []
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                entry = {"text": text, "level": para.level}
                entry.update(_para_style(para))
                paragraphs.append(entry)
    except Exception:
        pass
    return paragraphs


# ── SmartArt extraction ──────────────────────────────────────────────────────

def extract_smartart_text(shape):
    """
    SmartArt stores its data inside a <p:graphicFrame> element.
    We walk the XML looking for <a:t> text nodes that are children of
    SmartArt data namespaces.  Returns a list of plain strings.
    """
    texts = []
    try:
        from lxml import etree
        xml = shape._element
        # SmartArt text lives in dgm:pt/dgm:t or a:t nodes inside the graphic
        ns_map = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
        }
        # Collect all <a:t> nodes
        for t_node in xml.findall('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
            text = (t_node.text or '').strip()
            if text:
                texts.append(text)
    except Exception:
        pass
    return texts


def shape_is_smartart(shape):
    """Return True if a shape is a SmartArt graphic frame.

    Uses precise URI matching against the official DrawingML diagram namespace
    ('http://schemas.openxmlformats.org/drawingml/2006/diagram') to avoid
    false positives from shapes whose XML incidentally contains the words
    'diagram' or 'dgm' in comments or custom attributes.
    """
    _SMARTART_URI = 'http://schemas.openxmlformats.org/drawingml/2006/diagram'
    try:
        from lxml import etree
        tag = shape._element.tag
        if 'graphicFrame' not in tag:
            return False
        # Check graphicData uri attribute for the exact SmartArt namespace
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        for graphic_data in shape._element.iter(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData'
        ):
            uri = graphic_data.get('uri', '')
            if uri == _SMARTART_URI:
                return True
        return False
    except Exception:
        return False


# ── Table extraction ─────────────────────────────────────────────────────────

def extract_table(shape):
    """
    Extract a table shape into a structured dict:
    {
      "headers": ["col1", "col2", ...],
      "rows":    [["val", "val", ...], ...]
    }
    The first row is treated as headers if it contains any bold or
    visually distinct cells; otherwise an empty headers list is returned.
    """
    try:
        if not shape.has_table:
            return None
        table = shape.table
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text.strip())
            if any(cells):  # skip fully empty rows
                rows.append(cells)

        if not rows:
            return None

        # Heuristic: first row is headers if any cell is bold or ALL_CAPS
        def looks_like_header(row):
            return any(
                (c.isupper() and len(c) > 1) or c.endswith(':')
                for c in row if c
            )

        headers = []
        if rows and looks_like_header(rows[0]):
            headers = rows[0]
            rows = rows[1:]

        return {"headers": headers, "rows": rows}
    except Exception:
        return None


# ── Image extraction ─────────────────────────────────────────────────────────

def extract_image_from_shape(shape, slide_idx, img_idx, output_dir):
    """Extract and save image from a shape. Returns relative path or None."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_data = shape.image.blob
            img_ext = shape.image.ext or "png"
            img_filename = f"slide_{slide_idx + 1}_img_{img_idx + 1}.{img_ext}"
            img_path = output_dir / "images" / img_filename
            img_path.parent.mkdir(parents=True, exist_ok=True)
            with open(img_path, "wb") as f:
                f.write(img_data)
            return f"images/{img_filename}"
    except Exception as exc:
        import sys
        print(
            f"  ⚠  Slide {slide_idx + 1}, image {img_idx + 1}: "
            f"extraction failed — {type(exc).__name__}: {exc}",
            file=sys.stderr,
        )
    return None


# ── Layout detection ─────────────────────────────────────────────────────────

def detect_slide_layout(shapes, tables, smartarts):
    """
    Score-based layout detection with 8 layout types:
      title         — single dominant headline, minimal body
      section       — large centered text, transition/divider slide
      statement     — single bold statement, little else
      image-text    — image + text side-by-side or overlay
      content-list  — bulleted / numbered list
      table         — primary content is a table
      chart         — placeholder for chart slides (set externally)
      content       — generic fallback
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception:
        MSO_SHAPE_TYPE = None

    has_image = False
    has_table = bool(tables)
    has_smartart = bool(smartarts)
    text_blocks = []
    max_font_size = 0
    title_like = 0    # shapes whose name suggests title

    for shape in shapes:
        try:
            if MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
                continue
        except Exception:
            pass

        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()
            if not full_text:
                continue

            word_count = len(full_text.split())
            font_sizes = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)
                    except Exception:
                        pass

            top_size = max(font_sizes) if font_sizes else 0
            max_font_size = max(max_font_size, top_size)

            name_lower = shape.name.lower()
            is_title_named = any(k in name_lower for k in ("title", "heading"))
            if is_title_named:
                title_like += 1

            text_blocks.append({
                "words": word_count,
                "max_size": top_size,
                "title_named": is_title_named,
            })

    n_text = len(text_blocks)

    # ── Decision tree ──
    if has_table and n_text <= 2:
        return "table"

    if has_image and n_text >= 1:
        return "image-text"

    if n_text == 0 and has_smartart:
        return "content"  # SmartArt-only slide

    if n_text == 1:
        blk = text_blocks[0]
        if blk["max_size"] >= 32 or blk["title_named"]:
            # Single large text = either title or statement
            words = blk["words"]
            return "title" if words <= 12 else "statement"
        return "statement"

    if n_text == 2:
        sizes = [b["max_size"] for b in text_blocks]
        big, small = max(sizes), min(sizes)
        if big >= 28 and (big - small) >= 10:
            return "title"   # headline + subtitle pair
        return "content"

    # 3+ text shapes
    total_words = sum(b["words"] for b in text_blocks)
    if total_words >= 60 or n_text >= 4:
        return "content-list"

    return "content"


# ── Position-aware title detection ──────────────────────────────────────────

def _shape_top(shape):
    """Return shape top position in EMU (0 if unavailable)."""
    try:
        return shape.top or 0
    except Exception:
        return 0


def identify_title_shape(all_text_blocks_with_shapes):
    """
    Given a list of (shape, paragraphs) tuples, return the index most likely
    to be the slide title using a weighted scoring approach:
      - Shape name contains "title" (+3)
      - Shape is topmost on slide (+2)
      - First paragraph font size >= 24 (+1)
      - Single-paragraph shape (+1)
    """
    if not all_text_blocks_with_shapes:
        return -1

    scored = []
    min_top = min(_shape_top(s) for s, _ in all_text_blocks_with_shapes)

    for idx, (shape, paragraphs) in enumerate(all_text_blocks_with_shapes):
        score = 0
        name_lower = shape.name.lower()
        if "title" in name_lower:
            score += 3
        if _shape_top(shape) == min_top:
            score += 2
        if paragraphs and paragraphs[0].get("size", 0) >= 24:
            score += 1
        if len(paragraphs) == 1:
            score += 1
        scored.append((score, idx))

    scored.sort(key=lambda x: -x[0])
    return scored[0][1] if scored else 0


# ── Theme colour extraction ──────────────────────────────────────────────────

def _extract_theme_colors(prs) -> dict:
    """
    Extract the theme colour scheme from the first slide master.

    Returns a dict mapping role names to hex strings:
        {'dk1': '#1a1a1a', 'lt1': '#ffffff', 'accent1': '#4472c4', ...}

    Roles extracted: dk1, dk2, lt1, lt2, accent1–accent6.
    Returns an empty dict if extraction fails or no master exists.
    """
    _THEME_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    _SYSMAP   = {
        'windowText':   '#000000',
        'window':       '#ffffff',
        'background1':  '#ffffff',
        'text1':        '#000000',
        'background2':  '#eeece1',
        'text2':        '#1f497d',
    }

    try:
        masters = prs.slide_masters
        if not masters:
            return {}
        master = masters[0]
        theme_elem = master.element.find(
            './/{%s}theme' % _THEME_NS
        )
        if theme_elem is None:
            # Try via the theme part if the element isn't embedded
            try:
                from pptx.oxml.ns import qn
                theme_part = master.part.part_related_by(
                    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'
                )
                theme_elem = theme_part._element
            except Exception:
                pass

        if theme_elem is None:
            return {}

        # Walk <a:clrScheme> children
        clr_scheme = theme_elem.find(
            './/{%s}clrScheme' % _THEME_NS
        )
        if clr_scheme is None:
            return {}

        colors: dict[str, str] = {}
        for child in clr_scheme:
            role = child.tag.split('}')[-1]   # strip namespace
            # Each child has one sub-element: <a:srgbClr> or <a:sysClr>
            for sub in child:
                sub_tag = sub.tag.split('}')[-1]
                if sub_tag == 'srgbClr':
                    val = sub.get('val', '')
                    if val and len(val) == 6:
                        colors[role] = '#' + val.upper()
                elif sub_tag == 'sysClr':
                    # Prefer lastClr attribute (the resolved colour)
                    last = sub.get('lastClr', '')
                    if last and len(last) == 6:
                        colors[role] = '#' + last.upper()
                    else:
                        # Fall back to well-known system colour map
                        sys_name = sub.get('val', '')
                        fallback = _SYSMAP.get(sys_name)
                        if fallback:
                            colors[role] = fallback
                break  # only one sub-element expected per role

        return colors

    except Exception:
        return {}


# ── Main extraction ──────────────────────────────────────────────────────────

def extract_pptx(input_path: str, output_dir: str):
    """Main extraction function."""
    from pptx import Presentation

    input_path = Path(input_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    if not input_path.exists():
        print(f"❌ File not found: {input_path}")
        sys.exit(1)

    print(f"📂 Loading: {input_path.name}")

    try:
        prs = Presentation(str(input_path))
    except Exception as e:
        print(f"❌ Failed to open PPTX: {e}")
        sys.exit(1)

    # ── Extract theme colours from slide master ───────────────────────────────
    theme_colors = _extract_theme_colors(prs)
    if theme_colors:
        print(f"  🎨 Theme colours: {', '.join(f'{k}={v}' for k, v in list(theme_colors.items())[:4])}")

    slides_data = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "index": slide_idx,
            "number": slide_idx + 1,
            "layout": "content",
            "title": None,
            "subtitle": None,
            "body_paragraphs": [],
            "tables": [],
            "smartart": [],
            "images": [],
            "notes": None,
            "shape_count": len(slide.shapes),
        }

        # ── Extract notes ──
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data["notes"] = notes_text

        # ── Process shapes ──
        img_idx = 0
        text_shape_pairs = []   # (shape, paragraphs)

        for shape in slide.shapes:
            # Images
            img_path = extract_image_from_shape(shape, slide_idx, img_idx, output_dir)
            if img_path:
                slide_data["images"].append(img_path)
                img_idx += 1
                continue

            # Tables
            if shape.has_table:
                table_data = extract_table(shape)
                if table_data:
                    slide_data["tables"].append(table_data)
                continue

            # SmartArt
            if shape_is_smartart(shape):
                texts = extract_smartart_text(shape)
                if texts:
                    slide_data["smartart"].extend(texts)
                continue

            # Regular text
            paragraphs = extract_text_from_shape(shape)
            if paragraphs:
                text_shape_pairs.append((shape, paragraphs))

        # ── Identify title shape ──
        title_idx = identify_title_shape(text_shape_pairs)

        for i, (shape, paragraphs) in enumerate(text_shape_pairs):
            name_lower = shape.name.lower()

            if i == title_idx:
                slide_data["title"] = paragraphs[0]["text"]
                remaining = paragraphs[1:]
                if remaining:
                    slide_data["body_paragraphs"].extend(remaining)

            elif "subtitle" in name_lower and not slide_data["subtitle"]:
                slide_data["subtitle"] = paragraphs[0]["text"]
                remaining = paragraphs[1:]
                if remaining:
                    slide_data["body_paragraphs"].extend(remaining)

            else:
                slide_data["body_paragraphs"].extend(paragraphs)

        # ── Detect layout ──
        slide_data["layout"] = detect_slide_layout(
            slide.shapes,
            slide_data["tables"],
            slide_data["smartart"],
        )

        slides_data.append(slide_data)

        extras = []
        if slide_data["tables"]:
            extras.append(f"{len(slide_data['tables'])} table(s)")
        if slide_data["smartart"]:
            extras.append("SmartArt")
        if slide_data["images"]:
            extras.append(f"{len(slide_data['images'])} image(s)")
        extra_str = f"  [{', '.join(extras)}]" if extras else ""
        print(f"  ✓ Slide {slide_idx + 1}: [{slide_data['layout']}] "
              f"{slide_data['title'] or '(no title)'}{extra_str}")

    # ── Write slides.json ──
    json_path = output_dir / "slides.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump({
            "source_file": input_path.name,
            "total_slides": len(slides_data),
            "theme_colors": theme_colors,
            "slides": slides_data,
        }, f, ensure_ascii=False, indent=2)

    # ── Write summary.txt ──
    summary_path = output_dir / "summary.txt"
    with open(summary_path, "w", encoding="utf-8") as f:
        f.write("=== PPTX Content Summary ===\n")
        f.write(f"Source: {input_path.name}\n")
        f.write(f"Total slides: {len(slides_data)}\n\n")

        for slide in slides_data:
            f.write(f"--- Slide {slide['number']} [{slide['layout']}] ---\n")
            if slide["title"]:
                f.write(f"Title: {slide['title']}\n")
            if slide["subtitle"]:
                f.write(f"Subtitle: {slide['subtitle']}\n")
            if slide["body_paragraphs"]:
                f.write("Content:\n")
                for para in slide["body_paragraphs"]:
                    indent = "  " * para.get("level", 0)
                    style_tags = []
                    if para.get("bold"):
                        style_tags.append("bold")
                    if para.get("italic"):
                        style_tags.append("italic")
                    tag_str = f" [{', '.join(style_tags)}]" if style_tags else ""
                    f.write(f"  {indent}• {para['text']}{tag_str}\n")
            if slide["tables"]:
                for ti, tbl in enumerate(slide["tables"]):
                    f.write(f"Table {ti + 1}:\n")
                    if tbl["headers"]:
                        f.write(f"  Headers: {' | '.join(tbl['headers'])}\n")
                    for row in tbl["rows"][:5]:   # show first 5 rows in summary
                        f.write(f"  Row:     {' | '.join(row)}\n")
                    if len(tbl["rows"]) > 5:
                        f.write(f"  ... ({len(tbl['rows'])} rows total)\n")
            if slide["smartart"]:
                f.write(f"SmartArt: {' → '.join(slide['smartart'][:6])}\n")
            if slide["images"]:
                f.write(f"Images: {', '.join(slide['images'])}\n")
            if slide["notes"]:
                f.write(f"Notes: {slide['notes'][:120]}"
                        f"{'...' if len(slide['notes']) > 120 else ''}\n")
            f.write("\n")

    print(f"\n✅ Extraction complete!")
    print(f"   📄 Slides JSON : {json_path}")
    print(f"   📋 Summary     : {summary_path}")
    table_count = sum(len(s["tables"]) for s in slides_data)
    if table_count:
        print(f"   📊 Tables      : {table_count} extracted")
    smartart_count = sum(len(s["smartart"]) for s in slides_data)
    if smartart_count:
        print(f"   🔷 SmartArt    : {smartart_count} text items extracted")
    if any(s["images"] for s in slides_data):
        image_count = sum(len(s["images"]) for s in slides_data)
        print(f"   🖼️  Images      : {output_dir}/images/ ({image_count} files)")

    return slides_data


def main():
    parser = argparse.ArgumentParser(
        description="Extract content from PowerPoint files for frontend-slides skill (v2)"
    )
    parser.add_argument("input", help="Path to the .pptx file")
    parser.add_argument(
        "--output", "-o",
        default=".claude-design/pptx-extracted",
        help="Output directory (default: .claude-design/pptx-extracted)",
    )
    args = parser.parse_args()

    check_dependencies()
    extract_pptx(args.input, args.output)


if __name__ == "__main__":
    main()
